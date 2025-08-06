import os
import mimetypes
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from odf.opendocument import load

def scan_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    mime_type, _ = mimetypes.guess_type(file_path)

    suspicious_keywords = ['base64', 'vbscript', 'powershell', 'script', 'iframe', 'object', 'embed']
    result = []

    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                text = ''.join(page.extract_text() or '' for page in reader.pages)
                result = check_keywords(text, suspicious_keywords)
        elif ext == '.docx':
            doc = Document(file_path)
            text = ' '.join([p.text for p in doc.paragraphs])
            result = check_keywords(text, suspicious_keywords)
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text = ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            result = check_keywords(text, suspicious_keywords)
        elif ext == '.xlsx':
            wb = load_workbook(file_path, read_only=True)
            text = ' '.join([str(cell.value) for sheet in wb for row in sheet.iter_rows() for cell in row if cell.value])
            result = check_keywords(text, suspicious_keywords)
        elif ext in ['.odt', '.ods', '.odp']:
            text = str(load(file_path).text)
            result = check_keywords(text, suspicious_keywords)
        else:
            result.append('Unsupported file format.')

    except Exception as e:
        result.append(f'Error during scan: {str(e)}')

    return result if result else ['No threats detected.']

def check_keywords(text, keywords):
    found = [kw for kw in keywords if kw.lower() in text.lower()]
    return found if found else []
